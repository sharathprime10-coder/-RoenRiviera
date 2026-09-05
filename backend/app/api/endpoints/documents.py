import uuid
import io
import base64
import httpx
from fastapi import APIRouter, UploadFile, File, HTTPException, Depends
from langchain_community.vectorstores import SupabaseVectorStore
from langchain_core.documents import Document
from app.core.db import get_supabase_client, get_embeddings
from app.api.deps import verify_token
from app.core.config import settings

# Parsers
from pypdf import PdfReader
from pptx import Presentation
from docx import Document as DocxDocument

router = APIRouter()

@router.post("/upload")
async def upload_document(
    file: UploadFile = File(...),
    current_user: dict = Depends(verify_token)
):
    try:
        content = await file.read()
        filename = file.filename.lower() if file.filename else ""
        extension = filename.split('.')[-1] if '.' in filename else ''
        
        text_content = ""
        
        # 1. Parsing routing based on file extension
        if extension == 'pdf':
            reader = PdfReader(io.BytesIO(content))
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text_content += extracted + "\n\n"
                    
        elif extension in ['pptx', 'ppt']:
            prs = Presentation(io.BytesIO(content))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_content += shape.text + "\n"
                text_content += "\n"
                
        elif extension in ['docx', 'doc']:
            doc = DocxDocument(io.BytesIO(content))
            text_content = "\n".join([para.text for para in doc.paragraphs if para.text])
            
        elif extension in ['png', 'jpg', 'jpeg', 'webp']:
            # Primary: OpenRouter Vision API
            try:
                base64_image = base64.b64encode(content).decode('utf-8')
                headers = {
                    "Authorization": f"Bearer {settings.OPENROUTER_API_KEY}",
                    "HTTP-Referer": "http://localhost:8443",
                    "X-Title": "Campus Friend Vision Ingestion",
                }
                
                payload = {
                    # Using a strong vision model
                    "model": "google/gemini-flash-1.5-exp",
                    "messages": [
                        {
                            "role": "user",
                            "content": [
                                {
                                    "type": "text",
                                    "text": "Extract all readable text from this image exactly as it appears. Do not add conversational text, just the extracted text."
                                },
                                {
                                    "type": "image_url",
                                    "image_url": {
                                        "url": f"data:image/{extension};base64,{base64_image}"
                                    }
                                }
                            ]
                        }
                    ]
                }
                
                async with httpx.AsyncClient(timeout=30.0) as client:
                    response = await client.post("https://openrouter.ai/api/v1/chat/completions", headers=headers, json=payload)
                    response.raise_for_status()
                    result = response.json()
                    text_content = result['choices'][0]['message']['content']
                    
            except Exception as api_err:
                print(f"Vision API failed: {api_err}. Falling back to pytesseract...")
                # Secondary Fallback: pytesseract
                import pytesseract
                from PIL import Image
                img = Image.open(io.BytesIO(content))
                text_content = pytesseract.image_to_string(img)
        else:
            # Fallback for plain text files (.txt, .md, .csv)
            text_content = content.decode('utf-8')
            
        if not text_content.strip():
            raise HTTPException(status_code=400, detail=f"Could not extract any text from {file.filename}")

        # 2. Chunking strategy
        # Simple paragraph-based chunking strategy
        paragraphs = [p.strip() for p in text_content.split('\n\n') if p.strip()]
        
        docs = []
        doc_id = str(uuid.uuid4())
        
        for i, para in enumerate(paragraphs):
            # Only index chunks that have some length to avoid garbage tokens
            if len(para) > 10:
                docs.append(
                    Document(
                        page_content=para,
                        metadata={
                            "document_id": doc_id,
                            "document_name": file.filename,
                            "section": f"Chunk {i+1}"
                        }
                    )
                )
            
        if not docs:
            raise HTTPException(status_code=400, detail="Document contained no indexable text")

        # 3. Vector DB Upsert
        supabase = get_supabase_client()
        embeddings = get_embeddings()
        
        SupabaseVectorStore.from_documents(
            docs,
            embeddings,
            client=supabase,
            table_name="documents",
            query_name="match_documents"
        )
        
        return {"filename": file.filename, "status": "success", "chunks": len(docs)}
        
    except Exception as e:
        print(f"Error uploading document {file.filename}: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))
